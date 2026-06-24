#!/usr/bin/env python3
"""
Generate a standard Dukick Shooting Proposal deck (.pptx)
Usage: python shooting_proposal_generator.py --client "Brand Name" --project "Project Name"
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import argparse, sys

DARK   = RGBColor(26, 26, 26)
ACCENT = RGBColor(255, 77, 77)
WHITE  = RGBColor(255, 255, 255)
GRAY   = RGBColor(179, 179, 179)
BLACK  = RGBColor(40, 40, 40)

def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = DARK; bg.line.fill.background()
    # accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.5), Inches(1.5), Inches(0.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    # title
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title; p.font.size = Pt(54); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Arial"
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle; p2.font.size = Pt(22); p2.font.color.rgb = GRAY; p2.font.name = "Arial"; p2.space_before = Pt(14)
    return slide

def add_content_slide(prs, heading, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid(); header.fill.fore_color.rgb = DARK; header.line.fill.background()
    ht = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11.7), Inches(0.9))
    hp = ht.text_frame.paragraphs[0]
    hp.text = heading; hp.font.size = Pt(36); hp.font.bold = True; hp.font.color.rgb = WHITE; hp.font.name = "Arial"
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(1.5), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    ct = slide.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(11.5), Inches(5.2))
    tf = ct.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {b}"; p.font.size = Pt(20); p.font.color.rgb = BLACK; p.font.name = "Arial"; p.space_after = Pt(14)
    return slide

def build(client, project, out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, "SHOOTING PROPOSAL", f"Client: {client} – Dự án: {project}\nDukick PRODUCTION")
    add_content_slide(prs, "Tổng Quan Dự Án", [
        "Sản phẩm/Dịch vụ: [điền tại đây]",
        "Thời lượng đề xuất: 30–60 giây",
        "Định dạng đầu ra: Master 4K + Social cuts (9:16, 1:1, 16:9)",
        "Insight cốt lõi: [điền tại đây]",
        "Mục tiêu truyền thông: Awareness / Conversion / Loyalty",
        "Thời gian bàn giao dự kiến: [điền tại đây]",
    ])
    add_content_slide(prs, "Creative Direction", [
        "Tone \u0026 Mood: [Premium / Trẻ trung / Warm / Công nghệ / …]",
        "Màu sắc chủ đạo: [điền tại đây]",
        "Nhịp điệu (Pacing): [điền tại đây]",
        "Âm nhạc \u0026 SFX: [điền tại đây]",
        "Những điều KHÔNG làm (Don’ts): [điền tại đây]",
        "Key visual reference: [điền tại đây]",
    ])
    add_content_slide(prs, "Narrative \u0026 Storyboard", [
        "Mở đầu (0–5s): Hook – thu hút sự chú ý trong 3 giây đầu",
        "Thân bài (5–25s): Truyền tải thông điệp + lợi ích sản phẩm",
        "Kết (25–30s): Call-to-action + branding mạnh",
        "Số shot dự kiến: 18–22 shots",
        "Shot backup/dự phòng hậu kỳ: 3–5 shots",
        "Roll-out plan nếu cần pick-up: [điền tại đây]",
    ])
    add_content_slide(prs, "Kế Hoạch Sản Xuất", [
        "Pre-production: Concept final → Storyboard → Casting → Location recce → PPM",
        "Số ngày quay dự kiến: 1–2 ngày",
        "Số ngày Offline edit: 3–5 ngày",
        "Số ngày Online (grading + VFX + GFX): 2–3 ngày",
        "Địa điểm quay: Studio / Location / Green screen",
        "Casting: KOL / actor / chính chủ – cần profile trước PPM",
    ])
    add_content_slide(prs, "Ekip \u0026 Thiết Bị", [
        "Đạo diễn: [điền tại đây]",
        "DOP / Camera op: [điền tại đây]",
        "Producer / PM: [điền tại đây]",
        "Makeup \u0026 Hair + Stylist: [điền tại đây]",
        "Camera: [Model] + lens kit",
        "Lighting: [Kit] – thiết lập theo mood đã duyệt",
        "Gimbal / Drone / Motion control (nếu cần): [điền tại đây]",
    ])
    add_content_slide(prs, "Deliverables", [
        "Master TVC 4K (ProRes / H.264)",
        "Social cuts: 15s, 30s, 6s bumper",
        "Format ngang \u0026 dọc cho Reels / TikTok",
        "Frame grabs cho Key Visual / OOH",
        "Subtitles burned-in \u0026 file .srt riêng",
        "Moodboard / stills từ set nếu KH cần",
    ])
    add_content_slide(prs, "Timeline", [
        "Ngày 1–3: Finalize concept \u0026 storyboard",
        "Ngày 4–7: Casting, location booking, PPM",
        "Ngày 8–9: Quay (1–2 ngày)",
        "Ngày 10–14: Offline edit (client feedback round 1–2)",
        "Ngày 15–18: Online edit (grading, GFX, VFX)",
        "Ngày 19–20: Final delivery",
    ])
    add_title_slide(prs, "ĐẦU TƯ", "Báo giá chi tiết theo từng hạng mục – Dukick Production")
    add_content_slide(prs, "Báo Giá Tổng Quan", [
        "Pre-production: Concept + Storyboard + PPM",
        "Production: Crew + Equipment + Location + Talent + Logistics",
        "Post-production: Offline + Online + Sound + VO + VFX (nếu có)",
        "Deliverables: Master + Social cuts + Frame grabs",
        "Giá trên chưa bao gồm VAT (nếu phát sinh)",
        "Phát sinh ngoài scope sẽ được báo trước và phê duyệt bằng văn bản",
    ])
    prs.save(out_path)
    print(f"Saved: {out_path}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate Dukick Shooting Proposal pptx")
    parser.add_argument("--client", default="[Tên Khách Hàng]")
    parser.add_argument("--project", default="[Tên Dự Án]")
    parser.add_argument("--output", default="C:/DukickAgent/Dukick_Shooting_Proposal.pptx")
    args = parser.parse_args()
    build(args.client, args.project, args.output)
